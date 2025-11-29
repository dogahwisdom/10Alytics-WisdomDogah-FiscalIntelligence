"""
Convert Markdown Presentation to PowerPoint
Requires: python-pptx library
Install: pip install python-pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import re


def parse_markdown_slides(markdown_file: str) -> list:
    """
    Parse markdown file into slide content
    
    Args:
        markdown_file: Path to markdown file
    
    Returns:
        List of slide dictionaries
    """
    with open(markdown_file, 'r', encoding='utf-8') as f:
        content = f.read()
    
    # Split by slide separators
    slides = re.split(r'^---+\s*$', content, flags=re.MULTILINE)
    
    parsed_slides = []
    for slide in slides:
        slide = slide.strip()
        if not slide or slide.startswith('#'):
            continue
        
        lines = slide.split('\n')
        title = None
        content_lines = []
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
            
            # Extract title (first ## heading)
            if line.startswith('##') and not title:
                title = line.replace('##', '').strip()
            elif line.startswith('#'):
                # Other headings
                content_lines.append(line.replace('#', '').strip())
            else:
                content_lines.append(line)
        
        if title or content_lines:
            parsed_slides.append({
                'title': title or 'Slide',
                'content': content_lines
            })
    
    return parsed_slides


def create_powerpoint(markdown_file: str, output_file: str):
    """
    Create PowerPoint presentation from markdown
    
    Args:
        markdown_file: Path to markdown file
        output_file: Path to output PowerPoint file
    """
    slides = parse_markdown_slides(markdown_file)
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    for slide_data in slides:
        # Add slide
        slide_layout = prs.slide_layouts[1]  # Title and Content layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = slide_data['title']
        title_shape.text_frame.paragraphs[0].font.size = Pt(44)
        title_shape.text_frame.paragraphs[0].font.bold = True
        
        # Set content
        content_shape = slide.placeholders[1]
        tf = content_shape.text_frame
        tf.word_wrap = True
        
        for i, line in enumerate(slide_data['content']):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            
            p.text = line
            p.font.size = Pt(18)
            p.level = 0
            
            # Detect bullet points
            if line.startswith('-') or line.startswith('*'):
                p.level = 1
                p.text = line[1:].strip()
    
    # Save presentation
    prs.save(output_file)
    print(f"✓ PowerPoint presentation saved to: {output_file}")


if __name__ == "__main__":
    markdown_file = "presentation/PRESENTATION_SLIDES.md"
    output_file = "presentation/Fiscal_Intelligence_Presentation.pptx"
    
    try:
        create_powerpoint(markdown_file, output_file)
    except ImportError:
        print("Error: python-pptx not installed.")
        print("Install with: pip install python-pptx")
    except Exception as e:
        print(f"Error: {e}")

